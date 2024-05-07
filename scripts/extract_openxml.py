import os
import openpyxl
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

# Script variables
directory_current = os.getcwd()
directory_toolkit = os.path.dirname(directory_current)
default_input_directory = os.path.join(directory_toolkit, "_input", "bstrings")
default_output_directory = os.path.join(directory_toolkit, "_input", "bstrings")

def extract_text_from_docx(input_file, output_directory):
    # Open the DOCX file
    try:
        document = Document(input_file)
    except Exception as e:
        print("Error:", e)
        return False

    # Extract text from paragraphs
    text = []
    for paragraph in document.paragraphs:
        text.append(paragraph.text)

    # Generate output filename based on input filename
    output_filename = os.path.splitext(os.path.basename(input_file))[0] + "_docx.txt"
    output_path = os.path.join(output_directory, output_filename)

    # Write text to output file
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n\n'.join(text))
    except Exception as e:
        print("Error:", e)
        return False

    print("Text extracted successfully and saved to", output_path)
    
    # Remove the original DOCX file
    os.remove(input_file)
    print("Original DOCX file removed:", input_file)
    
    return True

def extract_text_from_xlsx(input_file, output_directory):
    # Open the XLSX file
    try:
        wb = openpyxl.load_workbook(input_file)
    except Exception as e:
        print("Error:", e)
        return False

    # Extract text from cells
    text = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell:
                    text.append(str(cell))

    # Generate output filename based on input filename
    output_filename = os.path.splitext(os.path.basename(input_file))[0] + "_xlsx.txt"
    output_path = os.path.join(output_directory, output_filename)

    # Write text to output file
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n\n'.join(text))
    except Exception as e:
        print("Error:", e)
        return False

    print("Text extracted successfully and saved to", output_path)
    
    # Remove the original XLSX file
    os.remove(input_file)
    print("Original XLSX file removed:", input_file)
    
    return True

def extract_text_from_pptx(input_file, output_directory):
    # Open the PPTX file
    try:
        presentation = Presentation(input_file)
    except Exception as e:
        print("Error:", e)
        return False

    # Extract text from slides
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    # Generate output filename based on input filename
    output_filename = os.path.splitext(os.path.basename(input_file))[0] + "_pptx.txt"
    output_path = os.path.join(output_directory, output_filename)

    # Write text to output file
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n\n'.join(text))
    except Exception as e:
        print("Error:", e)
        return False

    print("Text extracted successfully and saved to", output_path)
    
    # Remove the original PPTX file
    os.remove(input_file)
    print("Original PPTX file removed:", input_file)
    
    return True

def extract_text_from_pdf(input_file, output_directory):
    # Open the PDF file
    try:
        with open(input_file, 'rb') as f:
            reader = PdfReader(f)
            text = [page.extract_text() for page in reader.pages]
    except Exception as e:
        print("Error:", e)
        return False

    # Generate output filename based on input filename
    output_filename = os.path.splitext(os.path.basename(input_file))[0] + "_pdf.txt"
    output_path = os.path.join(output_directory, output_filename)

    # Write text to output file
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n\n'.join(text))
    except Exception as e:
        print("Error:", e)
        return False

    print("Text extracted successfully and saved to", output_path)
    
    # Remove the original PDF file
    os.remove(input_file)
    print("Original PDF file removed:", input_file)
    
    return True

def process_files(input_directory, output_directory):
    # Create output directory if it doesn't exist
    if not os.path.exists(output_directory):
        os.makedirs(output_directory)

    # Iterate over files in the input directory
    for filename in os.listdir(input_directory):
        input_file = os.path.join(input_directory, filename)
        if filename.endswith(".docx"):
            extract_text_from_docx(input_file, output_directory)
        elif filename.endswith(".xlsx"):
            extract_text_from_xlsx(input_file, output_directory)
        elif filename.endswith(".pptx"):
            extract_text_from_pptx(input_file, output_directory)
        elif filename.endswith(".pdf"):
            extract_text_from_pdf(input_file, output_directory)
        else:
            print("Skipping file:", input_file)

if __name__ == '__main__':
    process_files(default_input_directory, default_output_directory)
